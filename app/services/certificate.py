import subprocess
import smtplib
import time
import logging
from pathlib import Path
from email.message import EmailMessage
from datetime import datetime, timezone
from typing import Optional

from pptx import Presentation
from sqlmodel import Session

from app.core.config import settings
from app.db.session import engine
from app.db.schema import (
    EmailServiceJobStatus,
    EmailServiceRecipientStatus,
    EmailServiceTemplateType,
)
from app.services.database import DatabaseService
from app.core.exceptions import (
    JobNotFoundError,
    EventNotFoundError,
    TemplateNotFoundError,
    PdfConversionError,
    EmailSendError,
    CertificateGenerationError,
)

logger = logging.getLogger(__name__)


class CertificateService:
    def __init__(self):
        self.libreoffice = settings.get_libreoffice_path()

    def get_template_path(self, official: bool) -> str:
        if official:
            return settings.official_template
        return settings.unofficial_template

    def replace_placeholder(
        self,
        name: str,
        event_name: str,
        date: str,
        template_path: str,
        output_folder: Path,
    ) -> Path:
        template_file = Path(template_path)
        if not template_file.exists():
            raise TemplateNotFoundError(template_path)

        prs = Presentation(template_path)
        logger.info(f"Replacing placeholders in PPTX file: '{template_path}'")

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                    for run in paragraph.runs:
                        run_text: str = run.text
                        if (
                            settings.delimiter_start in run_text
                            and settings.delimiter_end in run_text
                        ):
                            placeholder = run_text[
                                run_text.find(settings.delimiter_start) : run_text.find(
                                    settings.delimiter_end
                                )
                                + len(settings.delimiter_end)
                            ]

                            name_placeholder = f"{settings.delimiter_start}{settings.name_placeholder}{settings.delimiter_end}"
                            event_placeholder = f"{settings.delimiter_start}{settings.event_name_placeholder}{settings.delimiter_end}"
                            date_placeholder = f"{settings.delimiter_start}{settings.date_placeholder}{settings.delimiter_end}"

                            if name_placeholder in run_text:
                                run.text = run.text.replace(placeholder, name)
                            if event_placeholder in run_text:
                                run.text = run.text.replace(placeholder, event_name)
                            if date_placeholder in run_text:
                                run.text = run.text.replace(placeholder, date)

        safe_name = self._sanitize_filename(name)
        output_pptx_name = f"{safe_name}-output-certificate.pptx"
        output_path = output_folder / output_pptx_name
        prs.save(str(output_path))

        assert output_path.exists(), (
            f"PPTX save succeeded but file missing: {output_path}"
        )

        logger.info(f"PPTX output saved to '{output_path}'")
        return output_path

    def pptx_to_pdf(self, input_pptx_path: Path, output_folder: Path) -> Path:
        logger.info(f"Converting to PDF with '{self.libreoffice}'")

        cmd = [
            self.libreoffice,
            "--headless",
            "--convert-to",
            settings.conversion_extension,
            str(input_pptx_path.name),
        ]

        logger.info(f"Running command: {' '.join(cmd)}")

        result = subprocess.run(
            cmd,
            cwd=str(output_folder),
            capture_output=True,
            text=True,
        )

        if result.returncode != 0:
            raise PdfConversionError(
                f"PDF conversion failed with return code {result.returncode}",
                details={
                    "return_code": result.returncode,
                    "stdout": result.stdout,
                    "stderr": result.stderr,
                    "input_file": str(input_pptx_path),
                },
            )

        pdf_name = input_pptx_path.stem + ".pdf"
        output_pdf = output_folder / pdf_name

        if not output_pdf.exists():
            raise PdfConversionError(
                f"LibreOffice returned success but PDF not found at {output_pdf}",
                details={"expected_path": str(output_pdf), "stdout": result.stdout},
            )

        logger.info(f"PDF output saved to '{output_pdf}'")
        return output_pdf

    def send_email(
        self,
        recipient: str,
        name: str,
        event_name: str,
        pdf_path: Path,
    ) -> None:
        subject = f"شهادة حضور {event_name}"

        try:
            with open(settings.email_template, "r", encoding="utf-8") as f:
                body = f.read()
            with open(pdf_path, "rb") as f:
                pdf_content = f.read()
        except FileNotFoundError as e:
            raise TemplateNotFoundError(
                settings.email_template,
                details={"context": "email_template", "error": str(e)},
            )

        msg = EmailMessage()
        msg["From"] = settings.sender_email
        msg["To"] = recipient
        msg["Subject"] = subject
        msg.set_content(
            "This email contains HTML. Please view it in an HTML-compatible client."
        )

        body = body.replace("[Name]", name)
        body = body.replace("[Event Name]", event_name)
        msg.add_alternative(body, subtype="html")

        msg.add_attachment(
            pdf_content,
            maintype="application",
            subtype="pdf",
            filename=f"{event_name} شهادة حضور.pdf",
        )

        logger.info(f"Sending email from {settings.sender_email} to {recipient}")

        last_error: Optional[Exception] = None
        for attempt in range(settings.max_retries):
            try:
                with smtplib.SMTP(
                    settings.smtp_host, settings.smtp_port, timeout=30
                ) as smtp:
                    smtp.starttls()
                    smtp.login(settings.sender_email, settings.app_password)
                    smtp.send_message(msg)
                    logger.info(f"Email sent successfully to {recipient}")
                    return
            except smtplib.SMTPException as e:
                last_error = e
                logger.warning(
                    f"SMTP error sending to {recipient} (attempt {attempt + 1}/{settings.max_retries}): {e}"
                )
            except Exception as e:
                last_error = e
                logger.warning(
                    f"Failed to send to {recipient} (attempt {attempt + 1}/{settings.max_retries}): {e}"
                )

            if attempt < settings.max_retries - 1:
                time.sleep(settings.email_delay)

        raise EmailSendError(
            recipient=recipient,
            reason=f"Failed after {settings.max_retries} attempts: {last_error}",
        )

    def _sanitize_filename(self, name: str) -> str:
        import re

        safe = re.sub(r'[<>:"/\\|?*]', "", name)
        safe = re.sub(r"\s+", "-", safe.strip())
        return safe

    def check_libreoffice(self) -> bool:
        try:
            result = subprocess.run(
                [self.libreoffice, "--version"],
                capture_output=True,
                text=True,
                timeout=10,
            )
            return result.returncode == 0
        except Exception:
            return False

    def check_smtp(self) -> bool:
        try:
            with smtplib.SMTP(
                settings.smtp_host, settings.smtp_port, timeout=10
            ) as smtp:
                smtp.starttls()
                smtp.login(settings.sender_email, settings.app_password)
                return True
        except Exception:
            return False


def process_certificate_attendance_job(job_id: str) -> None:
    service = CertificateService()

    with Session(engine) as session:
        db = DatabaseService(session)

        job = db.get_job(job_id)
        if not job:
            raise JobNotFoundError(job_id)

        event_id = job.event_id
        if not event_id:
            raise CertificateGenerationError(
                f"Job {job_id} has no event_id",
                details={"job_id": job_id},
            )

        event = db.get_event(event_id)
        if not event:
            raise EventNotFoundError(event_id)

        db.update_job_status(job_id, status=EmailServiceJobStatus.PROCESSING)

        template_type = (
            EmailServiceTemplateType.OFFICIAL
            if event.is_official
            else EmailServiceTemplateType.UNOFFICIAL
        )
        template_path = service.get_template_path(bool(event.is_official))

        event_date = event.start_datetime.strftime("%Y-%m-%d")

        recipients = db.get_recipients_for_job(job_id)

        for recipient in recipients:
            assert recipient.id, "Recipient ID should exist"
            if not recipient.member_id:
                db.update_recipient_status(
                    recipient.id,
                    EmailServiceRecipientStatus.FAILED,
                    error="No member_id associated",
                )
                continue

            member = db.get_member(recipient.member_id)
            if not member:
                logger.error(f"Member {recipient.member_id} not found")
                db.update_recipient_status(
                    recipient.id,
                    EmailServiceRecipientStatus.FAILED,
                    error="Member not found in database",
                )
                continue

            logger.info(f"Processing member {member.id}: {member.name}")

            pptx_path: Optional[Path] = None
            try:
                output_folder = (
                    Path(settings.certificates_folder)
                    / f"{event_id}-{db.sanitize_filename(event.name)}"
                )
                output_folder.mkdir(parents=True, exist_ok=True)

                pptx_path = service.replace_placeholder(
                    name=member.name,
                    event_name=event.name,
                    date=event_date,
                    template_path=template_path,
                    output_folder=output_folder,
                )

                pdf_path = service.pptx_to_pdf(pptx_path, output_folder)

                relative_path = (
                    f"{event_id}-{db.sanitize_filename(event.name)}/{pdf_path.name}"
                )
                db.create_certificate(
                    recipient_id=recipient.id,
                    certificate_path=relative_path,
                    template_type=template_type,
                )

                if not member.email:
                    db.update_recipient_status(
                        recipient.id,
                        EmailServiceRecipientStatus.FAILED,
                        error="Member has no email address",
                    )
                    continue

                service.send_email(
                    recipient=member.email,
                    name=member.name,
                    event_name=event.name,
                    pdf_path=pdf_path,
                )

                db.update_recipient_status(
                    recipient.id, EmailServiceRecipientStatus.SENT
                )

            except (PdfConversionError, EmailSendError, TemplateNotFoundError) as e:
                logger.error(f"Processing error for {member.name}: {e.message}")
                db.update_recipient_status(
                    recipient.id,
                    EmailServiceRecipientStatus.FAILED,
                    error=e.message,
                )
            except Exception as e:
                logger.exception(f"Unexpected error processing {member.name}: {e}")
                db.update_recipient_status(
                    recipient.id,
                    EmailServiceRecipientStatus.FAILED,
                    error=str(e),
                )
            finally:
                try:
                    if pptx_path:
                        pptx_path.unlink()
                except Exception as cleanup_error:
                    logger.warning(f"Failed to cleanup PPTX file: {cleanup_error}")

        progress = db.get_job_progress(job_id)

        if progress['failed'] == progress['total']:
            db.update_job_status(job_id, status=EmailServiceJobStatus.FAILED)
        else:
            db.update_job_status(job_id, status=EmailServiceJobStatus.COMPLETED)

        logger.info(f"Job {job_id} completed: {progress['successful']}/{progress['total']} successful")


def process_certificate_custom_job(job_id: str) -> None:
    service = CertificateService()

    with Session(engine) as session:
        db = DatabaseService(session)

        job = db.get_job(job_id)
        if not job:
            raise JobNotFoundError(job_id)

        event_name = job.event_name or "Unknown Event"
        event_date = (job.event_start_datetime or datetime.now(timezone.utc)).strftime("%Y-%m-%d")
        official = bool(job.is_official)

        db.update_job_status(job_id, status=EmailServiceJobStatus.PROCESSING)

        template_type = (
            EmailServiceTemplateType.OFFICIAL
            if official
            else EmailServiceTemplateType.UNOFFICIAL
        )
        template_path = service.get_template_path(official)

        recipients = db.get_recipients_for_job(job_id)

        for recipient in recipients:
            assert recipient.id, "Recipient ID should exist"

            name = recipient.name
            email = recipient.email

            if recipient.member_id:
                member = db.get_member(recipient.member_id)
                if not member:
                    logger.error(f"Member {recipient.member_id} not found")
                    db.update_recipient_status(
                        recipient.id,
                        EmailServiceRecipientStatus.FAILED,
                        error="Member not found in database",
                    )
                    continue
                name = member.name
                email = member.email

            if not name:
                name = "Unknown"

            if not email:
                logger.error(f"Recipient {recipient.id} has no email address")
                db.update_recipient_status(
                    recipient.id,
                    EmailServiceRecipientStatus.FAILED,
                    error="No email address",
                )
                continue

            logger.info(f"Processing recipient: {name} ({email})")

            pptx_path: Optional[Path] = None
            try:
                output_folder = (
                    Path(settings.certificates_folder) / f"custom-{job_id[:8]}"
                )
                output_folder.mkdir(parents=True, exist_ok=True)

                pptx_path = service.replace_placeholder(
                    name=name,
                    event_name=event_name,
                    date=event_date,
                    template_path=template_path,
                    output_folder=output_folder,
                )

                pdf_path = service.pptx_to_pdf(pptx_path, output_folder)

                relative_path = f"custom-{job_id[:8]}/{pdf_path.name}"
                db.create_certificate(
                    recipient_id=recipient.id,
                    certificate_path=relative_path,
                    template_type=template_type,
                )

                service.send_email(
                    recipient=email,
                    name=name,
                    event_name=event_name,
                    pdf_path=pdf_path,
                )

                db.update_recipient_status(
                    recipient.id, EmailServiceRecipientStatus.SENT
                )

            except (PdfConversionError, EmailSendError, TemplateNotFoundError) as e:
                logger.error(f"Processing error for {name}: {e.message}")
                db.update_recipient_status(
                    recipient.id, EmailServiceRecipientStatus.FAILED, error=e.message
                )
            except Exception as e:
                logger.exception(f"Unexpected error processing {name}: {e}")
                db.update_recipient_status(
                    recipient.id, EmailServiceRecipientStatus.FAILED, error=str(e)
                )
            finally:
                try:
                    if pptx_path:
                        pptx_path.unlink()
                except Exception as cleanup_error:
                    logger.warning(f"Failed to cleanup PPTX file: {cleanup_error}")

        progress = db.get_job_progress(job_id)

        if progress['failed'] == progress['total']:
            db.update_job_status(job_id, status=EmailServiceJobStatus.FAILED)
        else:
            db.update_job_status(job_id, status=EmailServiceJobStatus.COMPLETED)

        logger.info(f"Job {job_id} completed: {progress['successful']}/{progress['total']} successful")


certificate_service = CertificateService()
