import subprocess
import smtplib
import time
import logging
from pathlib import Path
from email.message import EmailMessage
from datetime import datetime
from typing import Optional

from pptx import Presentation

from config import settings
from models import JobStatus, MemberStatus, MemberResult, Member
from storage import storage

logger = logging.getLogger(__name__)


class CertificateService:
    """Service for generating and sending certificates."""

    def __init__(self):
        self.libreoffice = settings.get_libreoffice_path()

    def get_template_path(self, official: bool) -> str:
        """Get the certificate template path based on official flag."""
        if official:
            return settings.official_template
        return settings.unofficial_template

    def replace_placeholder(
        self,
        name: str,
        event_name: str,
        date: str,
        template_path: str,
        output_folder: Path,
    ) -> Path:
        """Replace placeholders in PPTX file and save to output folder."""
        prs = Presentation(template_path)
        logger.info(f"Replacing placeholders in PPTX file: '{template_path}'")

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run_text: str = run.text
                        if (
                            settings.delimiter_start in run_text
                            and settings.delimiter_end in run_text
                        ):
                            # Find placeholder
                            placeholder = run_text[
                                run_text.find(settings.delimiter_start) : run_text.find(
                                    settings.delimiter_end
                                )
                                + len(settings.delimiter_end)
                            ]

                            name_placeholder = f"{settings.delimiter_start}{settings.name_placeholder}{settings.delimiter_end}"
                            event_placeholder = f"{settings.delimiter_start}{settings.event_name_placeholder}{settings.delimiter_end}"
                            date_placeholder = f"{settings.delimiter_start}{settings.date_placeholder}{settings.delimiter_end}"

                            if name_placeholder in run_text:
                                run.text = run.text.replace(placeholder, name)
                            if event_placeholder in run_text:
                                run.text = run.text.replace(placeholder, event_name)
                            if date_placeholder in run_text:
                                run.text = run.text.replace(placeholder, date)

        # Create safe filename from name
        safe_name = self._sanitize_filename(name)
        output_pptx_name = f"{safe_name}-output-certificate.pptx"
        output_path = output_folder / output_pptx_name
        prs.save(str(output_path))
        logger.info(f"PPTX output saved to '{output_path}'")
        return output_path

    def pptx_to_pdf(self, input_pptx_path: Path, output_folder: Path) -> Optional[Path]:
        """Convert PPTX to PDF using LibreOffice."""
        logger.info(f"Converting to PDF with '{self.libreoffice}'")

        cmd = [
            self.libreoffice,
            "--headless",
            "--convert-to",
            settings.conversion_extension,
            str(input_pptx_path.name),
        ]

        logger.info(f"Running command: {' '.join(cmd)}")

        result = subprocess.run(
            cmd,
            cwd=str(output_folder),
            capture_output=True,
            text=True,
        )

        if result.returncode == 0:
            # LibreOffice outputs PDF with same name but .pdf extension
            pdf_name = input_pptx_path.stem + ".pdf"
            output_pdf = output_folder / pdf_name
            logger.info(f"PDF output saved to '{output_pdf}'")
            return output_pdf
        else:
            logger.error(
                f"PDF conversion failed with return code {result.returncode}\n"
                f"stdout: {result.stdout}\nstderr: {result.stderr}"
            )
            return None

    def send_email(
        self,
        recipient: str,
        name: str,
        event_name: str,
        announced_name: str,
        pdf_path: Path,
    ) -> tuple[bool, Optional[str]]:
        """Send certificate email to recipient. Returns (success, error_message)."""
        subject = f"شهادة حضور {event_name}"

        try:
            with open(settings.email_template, "r", encoding="utf-8") as f:
                body = f.read()
            with open(pdf_path, "rb") as f:
                pdf_content = f.read()
        except FileNotFoundError as e:
            error_msg = f"File not found: {e}"
            logger.error(error_msg)
            return False, error_msg

        # Create the email message
        msg = EmailMessage()
        msg["From"] = settings.sender_email
        msg["To"] = recipient
        msg["Subject"] = subject
        msg.set_content(
            "This email contains HTML. Please view it in an HTML-compatible client."
        )

        # Replace placeholders in HTML body
        body = body.replace("[Name]", name)
        body = body.replace("[Event Name]", event_name)
        body = body.replace("[Registered Name]", announced_name)
        msg.add_alternative(body, subtype="html")

        msg.add_attachment(
            pdf_content,
            maintype="application",
            subtype="pdf",
            filename=f"{event_name} شهادة حضور.pdf",
        )

        logger.info(f"Sending email from {settings.sender_email} to {recipient}")

        error_msg = "Unknown error"
        for attempt in range(settings.max_retries):
            try:
                with smtplib.SMTP(settings.smtp_host, settings.smtp_port) as smtp:
                    smtp.starttls()
                    smtp.login(settings.sender_email, settings.app_password)
                    smtp.send_message(msg)
                    logger.info(f"Email sent successfully to {recipient}")
                    return True, None
            except Exception as e:
                error_msg = str(e)
                logger.warning(
                    f"Failed to send to {recipient} (attempt {attempt + 1}/{settings.max_retries}): {error_msg}"
                )
                if attempt < settings.max_retries - 1:
                    time.sleep(settings.email_delay)

        return False, f"Failed after {settings.max_retries} attempts: {error_msg}"

    def _sanitize_filename(self, name: str) -> str:
        """Create a safe filename from a name."""
        # Replace spaces with hyphens, remove unsafe characters
        import re

        safe = re.sub(r'[<>:"/\\|?*]', "", name)
        safe = re.sub(r"\s+", "-", safe.strip())
        return safe

    def check_libreoffice(self) -> bool:
        """Check if LibreOffice is available."""
        try:
            result = subprocess.run(
                [self.libreoffice, "--version"],
                capture_output=True,
                text=True,
                timeout=10,
            )
            return result.returncode == 0
        except Exception:
            return False

    def check_smtp(self) -> bool:
        """Check if SMTP connection can be established."""
        try:
            with smtplib.SMTP(
                settings.smtp_host, settings.smtp_port, timeout=10
            ) as smtp:
                smtp.starttls()
                smtp.login(settings.sender_email, settings.app_password)
                return True
        except Exception:
            return False


def process_certificates_job(
    job_id: str,
    event_name: str,
    announced_name: str,
    date: str,
    official: bool,
    members: list[Member],
    folder_name: str,
) -> None:
    """Background task to process certificate generation and sending."""
    service = CertificateService()
    output_folder = storage.get_job_folder_path(folder_name)
    template_path = service.get_template_path(official)
    created_at = datetime.utcnow()

    # Update job status to processing
    storage.update_job_status(job_id, status=JobStatus.processing)

    member_results: list[MemberResult] = []

    for i, member in enumerate(members, 1):
        logger.info(f"Processing [{i}/{len(members)}]: {member.name}")

        result = MemberResult(
            name=member.name,
            email=member.email,
            gender=member.gender.value,
            status=MemberStatus.pending,
        )

        try:
            # Generate PPTX with placeholders replaced
            pptx_path = service.replace_placeholder(
                name=member.name,
                event_name=event_name,
                date=date,
                template_path=template_path,
                output_folder=output_folder,
            )

            # Convert to PDF
            pdf_path = service.pptx_to_pdf(pptx_path, output_folder)

            if pdf_path is None:
                result.status = MemberStatus.failed
                result.error = "PDF conversion failed"
                storage.update_job_status(
                    job_id,
                    increment_completed=True,
                    increment_failed=True,
                )
                member_results.append(result)
                continue

            # Send email
            success, error = service.send_email(
                recipient=member.email,
                name=member.name,
                event_name=event_name,
                announced_name=announced_name,
                pdf_path=pdf_path,
            )

            if success:
                result.status = MemberStatus.sent
                result.sent_at = datetime.utcnow()
                result.certificate_url = f"/certificates/{folder_name}/{pdf_path.name}"
                storage.update_job_status(
                    job_id,
                    increment_completed=True,
                    increment_successful=True,
                )
            else:
                result.status = MemberStatus.failed
                result.error = error
                result.certificate_url = f"/certificates/{folder_name}/{pdf_path.name}"
                storage.update_job_status(
                    job_id,
                    increment_completed=True,
                    increment_failed=True,
                )

        except Exception as e:
            logger.exception(f"Error processing {member.name}: {e}")
            result.status = MemberStatus.failed
            result.error = str(e)
            storage.update_job_status(
                job_id,
                increment_completed=True,
                increment_failed=True,
            )

        member_results.append(result)

        # Update summary.json after each member
        storage.write_summary(
            folder_name=folder_name,
            job_id=job_id,
            event_name=event_name,
            announced_name=announced_name,
            date=date,
            official=official,
            members=member_results,
            status=JobStatus.processing,
            created_at=created_at,
        )

    # Finalize job
    completed_at = datetime.utcnow()
    final_status = JobStatus.completed

    # Check if all failed
    if all(m.status == MemberStatus.failed for m in member_results):
        final_status = JobStatus.failed

    storage.update_job_status(job_id, status=final_status)
    storage.mark_event_completed(event_name)

    # Write final summary
    storage.write_summary(
        folder_name=folder_name,
        job_id=job_id,
        event_name=event_name,
        announced_name=announced_name,
        date=date,
        official=official,
        members=member_results,
        status=final_status,
        created_at=created_at,
        completed_at=completed_at,
    )

    logger.info(f"Job {job_id} completed with status: {final_status.value}")


# Global service instance
certificate_service = CertificateService()
