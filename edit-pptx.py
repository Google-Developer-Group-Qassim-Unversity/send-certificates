from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pprint import pprint
from platform import system
from pathlib import Path
import subprocess
import smtplib
from email.message import EmailMessage
from dotenv import load_dotenv
from os import getenv, makedirs
from sys import exit
import csv
from urllib import request
import io
import time
load_dotenv()

EVENT_NAME = "معسكر رؤية الحاسب من كاوست"
ANNOUNCED_EVENT_NAME = "معسكر رؤية الحاسب من كاوست"
DATE = "04/11/2025-06/11/2025"
DATA_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vSXh9S1AoUGgr2QAwJ1pBjBsY_USsZR2xyaZorT2auDfTMExNultaGtJeCXThHt7hxQaG726NIKF0sx/pub?gid=0&single=true&output=csv"

# PRESENTATION_FILE_NAME = "certificate.pptx"
PRESENTATION_FILE_NAME = "certificate unofficial.pptx"

MAX_RETRIES=3
EMAIL_DELAY=4
DELIMITER_START = "<<"
DELIMITER_END = ">>"
OUTPUT_PRESENATION_FILE_NAME = "output-certificate.pptx"
OUTPUT_FILE_NAME = "output-certificate"
CONVERSION_EXTENTION = "pdf"
TEXT_OUTPUT = "output.txt"
NAME_PLACEHOLDER = "name"
EVENT_NAME_PLACEHOLDER = "event_name"
DATE_PLACEHOLDER = "event_date"

OUTPUT_FOLDER = EVENT_NAME + "-" + "output-files"
SENDER_EMAIL = "gdg.qu1@gmail.com"

if system() == "Windows":
    libreoffice  = r"C:\Program Files\LibreOffice\program\soffice.exe"
elif system() == "Linux":
    libreoffice = "libreoffice"

def main():

    makedirs(OUTPUT_FOLDER, exist_ok=True)
    data = extract_data()

    print("adding some test recipients...")
    data.insert(0, ("Test", "gdg.qu1@gmail.com"))
    data.insert(len(data)//2, ("مرت على بالي", "albrrak337@gmail.com"))
    data.insert(len(data), ("عبدالاله عبدالعزيز منصور البراك", "albrrak773@gmail.com"))
    i = 1

    for name, email in data:
        print(f"-----------Sending [{i}/{len(data)}]-----------")
        output_file_name = replace_placeholder(name=name, output_prs_file_name=name + "-" + OUTPUT_PRESENATION_FILE_NAME)
        output_pdf = pptx_to_pdf(input_pptx_file_name=output_file_name, output_pdf_file_name=Path().joinpath(OUTPUT_FOLDER, name + "-" + OUTPUT_FILE_NAME))
        send_email(recipient=email, attachment_file_name=output_pdf, name=name)
        i += 1

def test(name: str, email: str | None = None):
    output_test_file = 'test-output-files'
    makedirs(output_test_file, exist_ok=True)
    output_file_name = replace_placeholder(name=name, output_prs_file_name=name + "-" + OUTPUT_PRESENATION_FILE_NAME, output_folder=output_test_file)
    output_pdf = pptx_to_pdf(input_pptx_file_name=output_file_name, output_pdf_file_name=Path().joinpath(output_test_file, name + "-" + OUTPUT_FILE_NAME), output_folder=output_test_file)
    if email:
        send_email(recipient=email, attachment_file_name=output_pdf, name=name)


def extract_data(url = DATA_URL):
    print("Extracting data from Google sheets...")
    for attempt in range(MAX_RETRIES):
        try:
            with request.urlopen(url) as response:
                csv_bytes = response.read()
                csv_text = csv_bytes.decode('utf-8')
                
            # Use csv.DictReader to parse as list of dicts
            csv_reader = csv.DictReader(io.StringIO(csv_text))

            data = [(row.get("name"), row.get("email")) for row in csv_reader]
            if "" in data:
                print(f"❌ Empty name or email found in the data.")
                exit(1)

            print(f"extracted \x1b[32m{len(data)}\x1b[0m records successfully ✅")
            return data
        except Exception as e:
            print(f"Failed to extract data from Google Sheets waiting for {EMAIL_DELAY} seconds... [{attempt+1}/{MAX_RETRIES}]")
            time.sleep(EMAIL_DELAY)
    print(f"❌ Failed to extract data from Google Sheets after {MAX_RETRIES} attempts.")
    exit(1)

def replace_placeholder(name,prs_file_name = PRESENTATION_FILE_NAME, output_prs_file_name = OUTPUT_PRESENATION_FILE_NAME, event_name = EVENT_NAME, date = DATE, output_folder = OUTPUT_FOLDER):
    prs = Presentation(prs_file_name)
    print(f"Replacing placeholders in PPTX file: \x1b[36m'{prs_file_name}'\x1b[0m...")

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run_text: str = run.text
                    if DELIMITER_START in run_text and DELIMITER_END in run_text:
                        placeholder = run_text[run_text.find(DELIMITER_START):run_text.find(DELIMITER_END) + len(DELIMITER_END)]
                        if DELIMITER_START + NAME_PLACEHOLDER + DELIMITER_END in run_text:
                            # print(f"Found PlaceHolder in \x1b[32m{run_text}\x1b[0m replacing {placeholder} with '{name}'")
                            run.text = run.text.replace(placeholder, name)
                        if DELIMITER_START + EVENT_NAME_PLACEHOLDER + DELIMITER_END in run_text:
                            # print(f"Found PlaceHolder in \x1b[32m{run_text}\x1b[0m replacing {placeholder} with '{event_name}'")
                            run.text = run.text.replace(placeholder, event_name)
                        if DELIMITER_START + DATE_PLACEHOLDER + DELIMITER_END in run_text:
                            run.text = run.text.replace(placeholder, date)
                    else:
                        # print(f"Skipping Run: \x1b[31m{run_text}\x1b[0m")
                        pass
    output = Path(".").joinpath(output_folder, output_prs_file_name)
    prs.save(output)
    print(f"PPTX output saved successfully ✅ to \x1b[36m'{output}'\x1b[0m")
    return output_prs_file_name

def pptx_to_pdf(input_pptx_file_name: str = OUTPUT_PRESENATION_FILE_NAME, output_pdf_file_name:str = OUTPUT_FILE_NAME, output_folder = OUTPUT_FOLDER):
    print(f"Saving As PDF with \x1b[36m'{libreoffice}'\x1b[0m")
    cmd = [
        libreoffice,
        "--headless",
        "--convert-to",
        CONVERSION_EXTENTION,
        str(input_pptx_file_name),
        str(output_pdf_file_name)
    ]
    print(f"Running command \x1b[36m'{" ".join(cmd)}'\x1b[0m...")
    result = subprocess.run(cmd, cwd=Path().joinpath(".", output_folder), capture_output=True, text=True)
    if result.returncode == 0:
        output = str(output_pdf_file_name) + "." + CONVERSION_EXTENTION
        print(f"command Ran successfully ✅")
        print(f"PDF Output saved to \x1b[36m'{output}'\x1b[0m")
        return output
    else:
        print(f"command Failed ❌ with return code \x1b[31m{result.returncode}\x1b[0m \n\x1b[33mstdout\x1b[0m:\n{result.stdout}\n\x1b[33mstderr:\x1b[0m\n{result.stderr}")
        exit(1)

def send_email(recipient: str = "albrrak773@gmail.com", html: str = "", attachment_file_name = OUTPUT_FILE_NAME, name = None,  event_name = EVENT_NAME, registered_name = ANNOUNCED_EVENT_NAME):
    subject = f"شهادة حضور {event_name}"

    if not name:
        system.exit("Name is required to send email.")
    if not recipient:
        system.exit("Recipient email is required to send email.")
    
    with open("index.html", 'r', encoding="utf-8") as f:
        body = f.read()
    with open(attachment_file_name, "rb") as f:
        pdf = f.read()

    # Create the email message
    msg = EmailMessage()
    msg["From"] = SENDER_EMAIL
    msg["To"] = recipient
    msg["Subject"] = subject
    msg.set_content("This email contains HTML. Please view it in an HTML-compatible client.")

    body = body.replace("[Name]", name)
    body = body.replace("[Event Name]", event_name)
    body = body.replace("[Registered Name]", registered_name)
    msg.add_alternative(body, subtype='html')
    
    msg.add_attachment(
        pdf,
        maintype="application",
        subtype="pdf",
        filename=f"{event_name} شهادة حضور" + ".pdf"
    )

    # Send the email via an SMTP server
    print(f"sending email from \x1b[32m{SENDER_EMAIL}\x1b[0m To \x1b[33m{recipient}\x1b[0m\nSubject \x1b[34m{subject}\x1b[0m with content \x1b[35m{body[:15]}\x1b[0m")

    for attempt in range(MAX_RETRIES):
        try:
            with smtplib.SMTP("smtp.gmail.com", 587) as smtp:
                print("Starting TLS...")
                smtp.starttls()
                print("Logging in...")
                smtp.login(SENDER_EMAIL, getenv("APP_PASSWORD"))
                print("Logged in Successfully ✅")
                print("Sending email...")
                smtp.send_message(msg)
                print(f"Email sent successfully ✅")
            break
        except Exception as e:
            print(f"Failed to send to {recipient} waiting for {EMAIL_DELAY}... [{attempt+1}/{MAX_RETRIES}]")
            time.sleep(EMAIL_DELAY)
    return

def extract_all_text_from_presentation(prs_file_name = PRESENTATION_FILE_NAME, output_txt_file_name = TEXT_OUTPUT ):
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    prs = Presentation(prs_file_name)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    with open(output_txt_file_name, 'w', encoding='utf-8') as f:
        for run in text_runs:
            f.write(run + "\n")


if __name__ == "__main__":
    # main()
    test('عبدالاله البراك')